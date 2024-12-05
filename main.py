import random
import string
from fastapi import FastAPI, File, UploadFile
from fastapi.middleware.cors import CORSMiddleware
from io import BytesIO
import fitz  # PyMuPDF
import pytesseract
from PIL import Image
import spacy
from docx import Document
from pptx import Presentation
import logging
import magic  # For MIME type detection

# Initialize FastAPI app
app = FastAPI()

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Load the NLP model for sentence parsing
nlp = spacy.load("en_core_web_sm")

# Configure Tesseract OCR if needed (adjust the path for your environment)
# pytesseract.pytesseract.tesseract_cmd = r"/path/to/tesseract"

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],  # Change to your frontend's origin in production
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Function to extract text from a PDF file and its embedded images
def extract_text_from_pdf_and_images(pdf_file: UploadFile):
    logger.info("Extracting text from PDF...")
    pdf_bytes = pdf_file.file.read()
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    text = ""
    for page in doc:
        text += page.get_text()
        for img in page.get_images(full=True):  # OCR for embedded images
            xref = img[0]
            base_image = doc.extract_image(xref)
            image_bytes = base_image["image"]
            image = Image.open(BytesIO(image_bytes))
            text += pytesseract.image_to_string(image)
    return text

# Function to extract text from a Word document
def extract_text_from_word(word_file: UploadFile):
    logger.info("Extracting text from Word document...")
    word_bytes = word_file.file.read()
    doc = Document(BytesIO(word_bytes))
    return "\n".join(para.text for para in doc.paragraphs)

# Function to extract text from a PowerPoint file
def extract_text_from_ppt(ppt_file: UploadFile):
    logger.info("Extracting text from PowerPoint...")
    ppt_bytes = ppt_file.file.read()
    prs = Presentation(BytesIO(ppt_bytes))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to determine file type using MIME type
def detect_file_type(file: UploadFile):
    logger.info("Detecting file type...")
    mime = magic.Magic(mime=True)
    mime_type = mime.from_buffer(file.file.read(1024))
    file.file.seek(0)  # Reset file pointer
    return mime_type

# Function to extract text from various file types
def extract_text_from_file(file: UploadFile):
    mime_type = detect_file_type(file)
    if "pdf" in mime_type:
        return extract_text_from_pdf_and_images(file)
    elif "word" in mime_type or "officedocument.wordprocessingml.document" in mime_type:
        return extract_text_from_word(file)
    elif "presentation" in mime_type or "officedocument.presentationml.presentation" in mime_type:
        return extract_text_from_ppt(file)
    else:
        raise ValueError("Unsupported file type. Please upload a PDF, Word, or PowerPoint file.")

# Function to identify sentences from extracted text
def identify_key_sentences(text):
    doc = nlp(text)
    return [sent.text.strip() for sent in doc.sents]

# Function to generate random multiple-choice questions
def generate_mcq(question, correct_answer, all_answers):
    choices = [correct_answer] + random.sample(all_answers, 3)  # Add 3 incorrect answers
    random.shuffle(choices)  # Shuffle the order
    return f"{question}\n" + "\n".join(f"{chr(65 + i)}) {choice}" for i, choice in enumerate(choices))

# Function to generate short-answer questions
def generate_short_answer(question):
    return f"{question}\n(Provide your answer here.)"

# Function to generate true/false questions
def generate_true_false(question, correct_answer):
    return f"{question}\n- True / False\nCorrect Answer: {correct_answer}"

# Function to analyze text and generate questions based on identified entities
def generate_questions_from_text(text):
    questions = []
    sentences = identify_key_sentences(text)
    for sentence in sentences:
        if random.random() < 0.33:  # MCQ
            question_text = f"What does the following mean: '{sentence}'?"
            questions.append(generate_mcq(question_text, sentence, ["Example A", "Example B", "Example C"]))
        elif random.random() < 0.5:  # Short Answer
            questions.append(generate_short_answer(f"Explain: '{sentence}'"))
        else:  # True/False
            questions.append(generate_true_false(f"Is the following correct? '{sentence}'", "True" if random.random() > 0.5 else "False"))
    return questions

@app.post("/generate_exam/")
async def generate_exam_from_file(file: UploadFile = File(...)):
    try:
        logger.info("Processing file: %s", file.filename)
        extracted_text = extract_text_from_file(file)
        questions = generate_questions_from_text(extracted_text)
        logger.info("Question generation complete.")
        return {"questions": questions}
    except ValueError as e:
        logger.error("Error: %s", e)
        return {"error": str(e)}
    except Exception as e:
        logger.error("Unexpected Error: %s", e)
        return {"error": "An unexpected error occurred."}
